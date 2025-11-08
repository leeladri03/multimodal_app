import streamlit as st
import pymysql
import fitz  # PyMuPDF
import io
from docx import Document
from pptx import Presentation
from PIL import Image
import google.generativeai as genai
import pandas as pd

# ----------------- MYSQL CONNECTION -----------------
DB_HOST = "localhost"
DB_USER = "root"
DB_PASS = "6309625235Ll@"  # your MySQL password
DB_NAME = "multimodal"

def get_connection():
    return pymysql.connect(host=DB_HOST, user=DB_USER, password=DB_PASS, database=DB_NAME)

# ----------------- GEMINI CONFIG -----------------
genai.configure(api_key="AIzaSyD5X-wZ6Nrrqe28WW_tnvFQtraJPMyop4s")
model = genai.GenerativeModel("models/gemini-2.5-flash")

# ----------------- FILE TEXT EXTRACTION -----------------
def extract_text(file, file_type):
    text = ""
    try:
        if file_type == "application/pdf":
            pdf = fitz.open(stream=file.read(), filetype="pdf")
            for page in pdf:
                text += page.get_text()
            pdf.close()

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(io.BytesIO(file.read()))
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif file_type.startswith("image/"):
            text = "[Image uploaded — will be analyzed by Gemini Vision when you ask questions.]"

        elif file_type.startswith("audio/"):
            text = "[Audio uploaded — Gemini can transcribe or summarize when queried.]"

        elif file_type.startswith("video/"):
            text = "[Video uploaded — Gemini will summarize or describe when queried.]"

        elif file_type == "text/plain":
            text = file.read().decode("utf-8")

        else:
            text = "[Unsupported file type — only PDF, DOCX, PPTX, images, audio, video, and text supported.]"

    except Exception as e:
        text = f"[Error extracting text: {e}]"

    return text

# ----------------- STREAMLIT PAGE CONFIG -----------------
st.set_page_config(page_title="Multimodal AI Dashboard", layout="wide", page_icon="🤖")

# ----------------- CUSTOM STYLING -----------------
st.markdown("""
    <style>
    .main {
        background-color: #f8fafc;
        padding: 20px;
    }
    .stApp {
        background-color: #f8fafc;
    }
    .title {
        background: linear-gradient(to right, #2563eb, #38bdf8);
        color: white;
        padding: 25px;
        border-radius: 12px;
        text-align: center;
        font-size: 30px;
        font-weight: bold;
        box-shadow: 0 4px 10px rgba(0,0,0,0.1);
        margin-bottom: 30px;
    }
    .card {
        background-color: white;
        border-radius: 12px;
        padding: 20px;
        box-shadow: 0 3px 10px rgba(0,0,0,0.05);
        margin-bottom: 20px;
    }
    .answer-box {
        background-color: #e0f2fe;
        padding: 20px;
        border-radius: 10px;
        margin-top: 10px;
    }
    </style>
""", unsafe_allow_html=True)

# ----------------- HEADER -----------------
st.markdown('<div class="title">🤖 Multimodal Data Processing Dashboard</div>', unsafe_allow_html=True)
st.write("Upload and analyze **PDF, DOCX, PPTX, TXT, images, audio, and video files** using Google Gemini AI and MySQL.")

# ----------------- LAYOUT -----------------
col1, col2 = st.columns([1, 2])

with col1:
    st.subheader("📂 Upload Files")
    uploaded_files = st.file_uploader(
        "Upload files (PDF, DOCX, PPTX, TXT, Images, Audio, Video):",
        accept_multiple_files=True,
        type=["pdf", "docx", "pptx", "txt", "jpg", "jpeg", "png", "mp3", "wav", "mp4", "mov"]
    )

file_data = []  # for Gemini input

if uploaded_files:
    for file in uploaded_files:
        file_type = file.type
        text = extract_text(file, file_type)

        try:
            conn = get_connection()
            cur = conn.cursor()
            sql = "INSERT INTO files (filename, filetype, content) VALUES (%s, %s, %s)"
            cur.execute(sql, (file.name, file_type, text))
            conn.commit()
            conn.close()
            st.success(f"✅ {file.name} saved to database.")
        except Exception as e:
            st.error(f"Database Error: {e}")

        st.markdown(f"#### 📄 {file.name}")
        st.text_area("Extracted Content Preview:", text[:700], height=150)

        if file_type.startswith("image/"):
            file.seek(0)
            image = Image.open(file)
            st.image(image, width=250)
            file_data.append(image)
        else:
            file_data.append(text)

# ----------------- QUERY SECTION -----------------
st.markdown("---")
st.subheader("💬 Ask Questions about Your Files")

query = st.text_input("Type your question here:")

if query and file_data:
    with st.spinner("Thinking... 🤔"):
        try:
            inputs = [query] + file_data
            response = model.generate_content(inputs)
            st.markdown("### 🤖 **Answer:**")
            st.markdown(f"<div class='answer-box'>{response.text}</div>", unsafe_allow_html=True)
        except Exception as e:
            st.error(f"Error: {e}")

# ----------------- FOOTER -----------------
st.markdown("""
---
<div style="text-align:center; color:gray; font-size:14px;">
Built with ❤️ using Streamlit, MySQL, and Google Gemini AI.
</div>
""", unsafe_allow_html=True)
